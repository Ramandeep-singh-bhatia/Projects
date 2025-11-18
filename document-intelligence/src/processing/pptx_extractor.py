"""
PPTX (PowerPoint) document extractor.
Handles PowerPoint presentations slide-by-slide.
"""

from pathlib import Path
from typing import Dict, Any
from pptx import Presentation
from src.processing.base_extractor import BaseExtractor, ExtractedDocument
from src.utils.logger import app_logger as logger


class PPTXExtractor(BaseExtractor):
    """Extractor for PPTX documents."""

    def can_handle(self, file_path: Path) -> bool:
        """Check if file is a PPTX."""
        return self._get_file_extension(file_path) == 'pptx'

    def extract(self, file_path: Path) -> ExtractedDocument:
        """
        Extract content and metadata from PPTX.

        Args:
            file_path: Path to PPTX file

        Returns:
            ExtractedDocument with content and metadata
        """
        logger.info(f"Extracting PPTX: {file_path}")

        try:
            prs = Presentation(file_path)

            # Extract content from all slides
            content = self._extract_content(prs)

            # Extract metadata
            metadata = self._extract_metadata(prs, file_path)

            word_count = self._count_words(content)
            page_count = len(prs.slides)

            logger.info(f"Successfully extracted PPTX: {page_count} slides")

            return ExtractedDocument(
                content=content,
                metadata=metadata,
                page_count=page_count,
                word_count=word_count,
                language="en",
                title=metadata.get('title'),
                author=metadata.get('author'),
                created_date=metadata.get('created_date'),
                modified_date=metadata.get('modified_date')
            )

        except Exception as e:
            logger.error(f"Failed to extract PPTX {file_path}: {str(e)}")
            raise

    def _extract_content(self, prs: Presentation) -> str:
        """Extract content from all slides."""
        content = []

        for slide_num, slide in enumerate(prs.slides, 1):
            content.append(f"\n=== Slide {slide_num} ===\n")

            # Extract text from all shapes
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)

                # Extract table content if present
                if shape.has_table:
                    table_text = self._extract_table(shape.table)
                    slide_text.append(f"\n[Table]\n{table_text}")

            content.append("\n".join(slide_text))

            # Extract notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text:
                    content.append(f"\n[Speaker Notes]\n{notes_text}")

        return "\n".join(content)

    def _extract_table(self, table) -> str:
        """Extract table content from a shape."""
        table_text = []

        for row in table.rows:
            row_text = "\t".join(cell.text for cell in row.cells)
            table_text.append(row_text)

        return "\n".join(table_text)

    def _extract_metadata(self, prs: Presentation, file_path: Path) -> Dict[str, Any]:
        """Extract PPTX metadata."""
        metadata = {}

        try:
            # Core properties
            core_props = prs.core_properties

            metadata['title'] = core_props.title or ""
            metadata['author'] = core_props.author or ""
            metadata['subject'] = core_props.subject or ""
            metadata['keywords'] = core_props.keywords or ""
            metadata['category'] = core_props.category or ""
            metadata['comments'] = core_props.comments or ""

            # Dates
            if core_props.created:
                metadata['created_date'] = core_props.created
            if core_props.modified:
                metadata['modified_date'] = core_props.modified

            # Presentation info
            metadata['slide_count'] = len(prs.slides)
            metadata['slide_width'] = prs.slide_width
            metadata['slide_height'] = prs.slide_height

        except Exception as e:
            logger.warning(f"Failed to extract PPTX metadata: {str(e)}")

        return metadata
